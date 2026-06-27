"""Load uploaded files into modality-aware RAG chunks."""

from __future__ import annotations

import hashlib
import io
import mimetypes
import uuid
from pathlib import Path
from typing import List

from langchain_text_splitters import RecursiveCharacterTextSplitter

from memory_agent.config import PDF_PAGES_PER_CHUNK
from memory_agent.rag.types import RagChunk

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

TEXT_SPLITTER = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    length_function=len,
    separators=["\n\n", "\n", ". ", "? ", "! ", "; ", ", ", " ", ""],
)

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg", ".opus"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".webm", ".avi", ".mkv"}
PDF_EXTENSIONS = {".pdf"}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".json", ".html", ".htm", ".xml",
    ".yaml", ".yml", ".py", ".js", ".ts", ".java", ".go", ".rs", ".sql",
}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


def _chunk_id(source: str, index: int, suffix: str = "") -> str:
    digest = hashlib.sha256(f"{source}:{index}:{suffix}".encode("utf-8")).hexdigest()[:16]
    return digest


def _guess_mime(filename: str, fallback: str = "application/octet-stream") -> str:
    mime, _ = mimetypes.guess_type(filename)
    return mime or fallback


def _extension(filename: str) -> str:
    return Path(filename).suffix.lower()


def _base_metadata(source: str, doc_id: str, chunk_index: int, modality: str) -> dict:
    return {
        "source": source,
        "doc_id": doc_id,
        "chunk_index": chunk_index,
        "modality": modality,
    }


def _load_pdf_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(io.BytesIO(raw_bytes))
    total_pages = len(reader.pages)
    chunks: List[RagChunk] = []

    for batch_start in range(0, total_pages, PDF_PAGES_PER_CHUNK):
        writer = PdfWriter()
        batch_end = min(batch_start + PDF_PAGES_PER_CHUNK, total_pages)
        for page in reader.pages[batch_start:batch_end]:
            writer.add_page(page)

        buffer = io.BytesIO()
        writer.write(buffer)
        pdf_bytes = buffer.getvalue()
        chunk_index = batch_start // PDF_PAGES_PER_CHUNK
        chunk_id = _chunk_id(filename, chunk_index, "pdf")

        chunks.append(
            RagChunk(
                chunk_id=chunk_id,
                source=filename,
                modality="pdf",
                mime_type="application/pdf",
                text=f"PDF pages {batch_start + 1}-{batch_end} of {total_pages}",
                media_bytes=pdf_bytes,
                metadata={
                    **_base_metadata(filename, doc_id, chunk_index, "pdf"),
                    "page_start": batch_start + 1,
                    "page_end": batch_end,
                    "total_pages": total_pages,
                },
            )
        )
    return chunks


def _load_text_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    text = raw_bytes.decode("utf-8", errors="ignore")
    if not text.strip():
        return []

    segments = TEXT_SPLITTER.split_text(text)
    chunks: List[RagChunk] = []
    for index, segment in enumerate(segments):
        chunks.append(
            RagChunk(
                chunk_id=_chunk_id(filename, index, "text"),
                source=filename,
                modality="text",
                mime_type=_guess_mime(filename, "text/plain"),
                text=segment,
                metadata=_base_metadata(filename, doc_id, index, "text"),
            )
        )
    return chunks


def _load_docx_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    from docx import Document as DocxDocument

    document = DocxDocument(io.BytesIO(raw_bytes))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    return _load_text_chunks(filename, text.encode("utf-8"), doc_id)


def _load_xlsx_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    lines: List[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                lines.append(" | ".join(values))
    text = "\n".join(lines)
    return _load_text_chunks(filename, text.encode("utf-8"), doc_id)


def _load_pptx_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(raw_bytes))
    lines: List[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines.append(f"# Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    text = "\n".join(lines)
    return _load_text_chunks(filename, text.encode("utf-8"), doc_id)


def _load_image_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    ext = _extension(filename)
    mime = "image/jpeg" if ext in {".jpg", ".jpeg"} else _guess_mime(filename, "image/png")
    if ext not in IMAGE_EXTENSIONS:
        raise ValueError(f"Unsupported image type: {filename}")

    return [
        RagChunk(
            chunk_id=_chunk_id(filename, 0, "image"),
            source=filename,
            modality="image",
            mime_type=mime,
            text=f"Image file: {filename}",
            media_bytes=raw_bytes,
            metadata=_base_metadata(filename, doc_id, 0, "image"),
        )
    ]


def _load_audio_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    mime = _guess_mime(filename, "audio/mpeg")
    return [
        RagChunk(
            chunk_id=_chunk_id(filename, 0, "audio"),
            source=filename,
            modality="audio",
            mime_type=mime,
            text=f"Audio file: {filename}",
            media_bytes=raw_bytes,
            metadata={
                **_base_metadata(filename, doc_id, 0, "audio"),
                "note": "Embedded natively; max ~180s per Gemini Embedding 2 request.",
            },
        )
    ]


def _load_video_chunks(filename: str, raw_bytes: bytes, doc_id: str) -> List[RagChunk]:
    mime = _guess_mime(filename, "video/mp4")
    return [
        RagChunk(
            chunk_id=_chunk_id(filename, 0, "video"),
            source=filename,
            modality="video",
            mime_type=mime,
            text=f"Video file: {filename}",
            media_bytes=raw_bytes,
            metadata={
                **_base_metadata(filename, doc_id, 0, "video"),
                "note": "Embedded natively; max ~120s per Gemini Embedding 2 request.",
            },
        )
    ]


def load_uploaded_file(filename: str, raw_bytes: bytes) -> List[RagChunk]:
    """Load any supported upload into RAG chunks for indexing."""
    if not raw_bytes:
        return []

    doc_id = str(uuid.uuid4())
    ext = _extension(filename)

    if ext in PDF_EXTENSIONS:
        return _load_pdf_chunks(filename, raw_bytes, doc_id)
    if ext in IMAGE_EXTENSIONS:
        return _load_image_chunks(filename, raw_bytes, doc_id)
    if ext in AUDIO_EXTENSIONS:
        return _load_audio_chunks(filename, raw_bytes, doc_id)
    if ext in VIDEO_EXTENSIONS:
        return _load_video_chunks(filename, raw_bytes, doc_id)
    if ext in OFFICE_EXTENSIONS:
        if ext == ".docx":
            return _load_docx_chunks(filename, raw_bytes, doc_id)
        if ext == ".xlsx":
            return _load_xlsx_chunks(filename, raw_bytes, doc_id)
        if ext == ".pptx":
            return _load_pptx_chunks(filename, raw_bytes, doc_id)
    if ext in TEXT_EXTENSIONS or ext == "":
        return _load_text_chunks(filename, raw_bytes, doc_id)

    # Fallback: attempt UTF-8 text extraction before rejecting.
    try:
        text = raw_bytes.decode("utf-8")
        if text.strip():
            return _load_text_chunks(filename, raw_bytes, doc_id)
    except UnicodeDecodeError:
        pass

    raise ValueError(
        f"Unsupported file type: {filename}. "
        "Supported: text, PDF, images (PNG/JPEG/WebP/GIF), audio (MP3/WAV), "
        "video (MP4/MOV), and Office files (DOCX/XLSX/PPTX)."
    )
